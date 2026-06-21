# -*- coding: utf-8 -*-
"""PPTX export (python-pptx): an executive deck covering research + GTM strategy
+ marketing content. Dark title/closing slides, light content slides, consistent
palette, no accent stripes."""
import os
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x21, 0x29, 0x34)
MUTE = RGBColor(0x5A, 0x64, 0x72)
W, H = Inches(13.333), Inches(7.5)


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    first = True
    for text, size, color, bold in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Calibri"
    return tb


def _bullets(slide, x, y, w, h, items, size=14, color=INK, bullet="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"{bullet} {it}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def _title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    _text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.5),
          [("MARKET INTELLIGENCE & GTM", 16, ICE, True)])
    _text(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(2.0),
          [(title, 44, WHITE, True)])
    _text(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.0),
          [(subtitle, 20, ICE, False)])


def _header(s, kicker, title):
    _bg(s, WHITE)
    _text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
          [(kicker.upper(), 13, MUTE, True)])
    _text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
          [(title, 32, NAVY, True)])


def _card(s, x, y, w, h, fill):
    sh = s.shapes.add_shape(1, x, y, w, h)  # rectangle
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _content_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def export_pptx(result: Dict[str, Any], path: Optional[str] = None) -> str:
    r = result.get("report", {}) or {}
    plan = result.get("plan", {}) or {}
    gtm = result.get("gtm_strategy") or {}
    content = result.get("content") or {}
    name = plan.get("subject_entity") or r.get("title") or "Market Analysis"

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1. Title
    _title_slide(prs, name, "Executive Market Brief, GTM Strategy & Content")

    # 2. Executive summary
    s = _content_slide(prs)
    _header(s, "Executive Summary", name)
    _text(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(3.6),
          [(r.get("executive_summary", "—"), 16, INK, False)])
    conf = r.get("confidence_level", "—")
    _text(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.6),
          [(f"Confidence: {conf}", 13, MUTE, True)])
    if r.get("_evidence_limitation"):
        _text(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.7),
              [(r["_evidence_limitation"], 12, RGBColor(0xB0, 0x45, 0x00), True)])

    # 3. SWOT 2x2
    sw = r.get("subject_swot") or {}
    if any(sw.get(k) for k in ("strengths", "weaknesses", "opportunities", "threats")):
        s = _content_slide(prs)
        _header(s, "Position", "SWOT")
        quad = [("Strengths", "strengths", ICE), ("Weaknesses", "weaknesses", RGBColor(0xF3, 0xD9, 0xD9)),
                ("Opportunities", "opportunities", RGBColor(0xD9, 0xF0, 0xE1)),
                ("Threats", "threats", RGBColor(0xF6, 0xE7, 0xC8))]
        xs = [Inches(0.7), Inches(6.95)]
        ys = [Inches(1.9), Inches(4.6)]
        for i, (label, key, col) in enumerate(quad):
            x, y = xs[i % 2], ys[i // 2]
            _card(s, x, y, Inches(5.65), Inches(2.5), col)
            _text(s, x + Inches(0.2), y + Inches(0.12), Inches(5.3), Inches(0.4),
                  [(label, 16, NAVY, True)])
            pts = [p.get("point", "") if isinstance(p, dict) else str(p)
                   for p in (sw.get(key) or [])[:3]]
            _bullets(s, x + Inches(0.2), y + Inches(0.6), Inches(5.25), Inches(1.8), pts or ["—"], size=12)

    # 4. Competitors
    comps = r.get("company_competitors") or r.get("product_competitors") or []
    if comps:
        s = _content_slide(prs)
        _header(s, "Landscape", "Key Competitors")
        rows = [["Competitor", "Positioning / Note"]]
        for c in comps[:6]:
            note = c.get("value_proposition") or c.get("positioning") or c.get("directness") or ""
            rows.append([c.get("name", "—"), str(note)[:90]])
        tbl = s.shapes.add_table(len(rows), 2, Inches(0.7), Inches(1.9),
                                 Inches(11.9), Inches(0.4) * len(rows)).table
        tbl.columns[0].width = Inches(3.3)
        tbl.columns[1].width = Inches(8.6)
        for ci, head in enumerate(rows[0]):
            cell = tbl.cell(0, ci)
            cell.text = head
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        for ri in range(1, len(rows)):
            for ci in range(2):
                cell = tbl.cell(ri, ci)
                cell.text = rows[ri][ci]
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = INK

    # 5. GTM strategy
    if gtm:
        _f = gtm.get("foundation", {}) or {}
        _a = gtm.get("activation", {}) or {}
        _x = gtm.get("execution", {}) or {}
        s = _content_slide(prs)
        _header(s, "Go-To-Market", "Strategy")
        _text(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(0.9),
              [("North Star: ", 16, NAVY, True),
               (gtm.get("north_star") or _f.get("north_star", "—"), 16, INK, False)])
        _text(s, Inches(0.7), Inches(2.75), Inches(11.9), Inches(1.1),
              [("Positioning. ", 15, NAVY, True),
               (_f.get("positioning_statement", "—"), 14, INK, False)])
        _text(s, Inches(0.7), Inches(3.95), Inches(5.6), Inches(0.4),
              [("Competitive Wedge", 15, NAVY, True)])
        wedge = [d.get("sharpest_message", "") for d in
                 (_f.get("competitive_differentiation") or [])[:4]]
        _bullets(s, Inches(0.7), Inches(4.4), Inches(5.6), Inches(2.4), wedge or ["—"], size=12)
        _text(s, Inches(6.7), Inches(3.95), Inches(5.9), Inches(0.4),
              [(f"Motion: {(_a.get('motion') or {}).get('primary', '—')}", 15, NAVY, True)])
        roadmap = [f"{p.get('phase', '')}: {p.get('objective', '')}"
                   for p in (_x.get("roadmap_90day") or [])[:4]]
        _bullets(s, Inches(6.7), Inches(4.4), Inches(5.9), Inches(2.4),
                 roadmap or ["—"], size=12, bullet="→")


    # 6. Content highlights
    if content:
        s = _content_slide(prs)
        _header(s, f"Marketing Content (Phase {content.get('phase', 'A')})", "Launch-Ready Assets")
        cols = [("LinkedIn", [p.get("hook", "") for p in content.get("linkedin_posts", [])]),
                ("Blog", [b.get("title", "") for b in content.get("blog_drafts", [])]),
                ("Email", [e.get("subject", "") for e in content.get("email_drafts", [])])]
        x = Inches(0.7)
        for label, items in cols:
            _card(s, x, Inches(1.9), Inches(3.9), Inches(4.6), RGBColor(0xF2, 0xF5, 0xFC))
            _text(s, x + Inches(0.25), Inches(2.05), Inches(3.4), Inches(0.4),
                  [(label, 17, NAVY, True)])
            _bullets(s, x + Inches(0.25), Inches(2.6), Inches(3.45), Inches(3.7),
                     [i for i in items if i][:4] or ["—"], size=12)
            x += Inches(4.05)

    # 7. Recommendations (dark closing)
    s = _content_slide(prs)
    _bg(s, NAVY)
    _text(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.8),
          [("Recommendations", 32, WHITE, True)])
    recs = r.get("recommendations") or (gtm.get("quick_wins") if gtm else []) or ["—"]
    _bullets(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(4.8),
             [str(x) for x in recs[:6]], size=18, color=ICE)

    out = path or os.path.join("exports", "pptx",
                               f"gtm_{name.lower().replace(' ', '_')[:40]}.pptx")
    os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
    prs.save(out)
    print(f"   ✓ PPTX written: {out}")
    return out
